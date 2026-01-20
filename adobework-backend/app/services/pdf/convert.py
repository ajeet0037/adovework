"""
Advanced PDF Conversion Services
High-quality conversions with font and formatting preservation
Uses PyMuPDF (fitz) for professional-grade results
Includes OCR support for image-based PDFs
"""
import os
import io
from pathlib import Path
from typing import Optional, List
import fitz  # PyMuPDF
from pdf2docx import Converter as PDFToDocxConverter
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pdf2image import convert_from_path
from docx import Document as DocxDocument
from docx.shared import Pt as DocxPt, Inches as DocxInches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from PyPDF2 import PdfReader
import pytesseract
from PIL import Image
import cv2
import numpy as np

# camelot is optional - requires ghostscript
try:
    import camelot
    CAMELOT_AVAILABLE = True
except ImportError:
    CAMELOT_AVAILABLE = False


def analyze_pdf_type(pdf_path: str) -> dict:
    """
    Analyze PDF to determine the best conversion strategy.
    
    Returns dict with:
    - is_scanned: True if PDF is purely scanned/image-based (needs OCR)
    - is_complex: True if PDF has complex layout with images (like Aadhaar, ID cards)
    - has_text: True if PDF has extractable text
    - image_count: Count of images in the PDF
    - text_length: Total text length
    - recommendation: 'ocr', 'text', 'hybrid', or 'image'
    """
    try:
        pdf_doc = fitz.open(pdf_path)
        total_text_length = 0
        total_images = 0
        page_count = len(pdf_doc)
        
        # Check all pages (limit to 10 for performance)
        for page_num in range(min(10, page_count)):
            page = pdf_doc[page_num]
            text = page.get_text()
            total_text_length += len(text.strip())
            total_images += len(page.get_images())
        
        pdf_doc.close()
        
        # Analyze the PDF
        has_text = total_text_length > 100
        has_many_images = total_images >= page_count  # At least 1 image per page
        is_scanned = total_text_length < 100 and total_images > 0
        
        # Complex PDFs have both text AND significant images (like Aadhaar)
        is_complex = has_text and has_many_images
        
        # Determine recommendation
        if is_scanned:
            recommendation = 'ocr'  # Need OCR for scanned docs
        elif is_complex:
            recommendation = 'hybrid'  # Complex layout - use image mode with OCR text
        elif has_text and not has_many_images:
            recommendation = 'text'  # Simple text PDF - extract editable text
        else:
            recommendation = 'image'  # Default to image mode for safety
        
        return {
            'is_scanned': is_scanned,
            'is_complex': is_complex,
            'has_text': has_text,
            'image_count': total_images,
            'text_length': total_text_length,
            'page_count': page_count,
            'recommendation': recommendation
        }
    except Exception as e:
        print(f"PDF analysis error: {e}")
        return {
            'is_scanned': False,
            'is_complex': True,
            'has_text': False,
            'image_count': 0,
            'text_length': 0,
            'page_count': 0,
            'recommendation': 'image'  # Default to safe option
        }


def is_image_based_pdf(pdf_path: str) -> bool:
    """
    Check if PDF is image-based (scanned) or text-based
    Returns True if PDF has minimal text and likely needs OCR
    """
    analysis = analyze_pdf_type(pdf_path)
    return analysis['is_scanned']


def ocr_pdf_page(page, lang='eng'):
    """
    Perform OCR on a PDF page
    Returns extracted text with basic formatting
    """
    try:
        # Render page as high-quality image
        zoom = 2.0  # 2x zoom for better OCR accuracy
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        
        # Convert to PIL Image
        img_data = pix.tobytes("png")
        img = Image.open(io.BytesIO(img_data))
        
        # Convert to OpenCV format for preprocessing
        img_cv = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
        
        # Preprocessing for better OCR
        # Convert to grayscale
        gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
        
        # Apply thresholding to get black text on white background
        _, thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        # Denoise
        denoised = cv2.fastNlMeansDenoising(thresh, None, 10, 7, 21)
        
        # Convert back to PIL Image
        processed_img = Image.fromarray(denoised)
        
        # Perform OCR with better config
        custom_config = r'--oem 3 --psm 6'
        text = pytesseract.image_to_string(processed_img, lang=lang, config=custom_config)
        
        return text
    except Exception as e:
        print(f"OCR error: {e}")
        return ""


def pdf_to_word_image_mode(pdf_path: str, output_path: Optional[str] = None, dpi: int = 200) -> str:
    """
    Convert PDF to Word by rendering each page as a high-quality image.
    Best for complex layouts like ID cards, Aadhaar, certificates, etc.
    Preserves visual layout perfectly.
    
    Args:
        pdf_path: Path to input PDF
        output_path: Optional output path
        dpi: Image resolution (default 200 for good quality/size balance)
    
    Returns:
        Path to output Word document
    """
    if output_path is None:
        output_path = str(Path(pdf_path).with_suffix('.docx'))
    
    # Open PDF
    pdf_doc = fitz.open(pdf_path)
    doc = DocxDocument()
    
    # Calculate zoom factor for desired DPI (default PDF is 72 DPI)
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    
    for page_num in range(len(pdf_doc)):
        page = pdf_doc[page_num]
        
        # Add page break between pages (except first)
        if page_num > 0:
            doc.add_page_break()
        
        # Render page as high-quality image
        pix = page.get_pixmap(matrix=mat)
        
        # Save temp image
        temp_img = f"/tmp/pdf_page_{page_num}_{os.getpid()}.png"
        pix.save(temp_img)
        
        # Get image dimensions to maintain aspect ratio
        img_width_inches = page.rect.width / 72.0  # Convert points to inches
        img_height_inches = page.rect.height / 72.0
        
        # Scale to fit within Word page (max 6.5" wide for standard margins)
        max_width = 6.5
        if img_width_inches > max_width:
            scale = max_width / img_width_inches
            img_width_inches = max_width
            img_height_inches *= scale
        
        # Insert image into Word document
        try:
            doc.add_picture(temp_img, width=DocxInches(img_width_inches))
        except Exception as e:
            print(f"Error adding image for page {page_num}: {e}")
        
        # Cleanup temp file
        try:
            os.remove(temp_img)
        except:
            pass
    
    pdf_doc.close()
    doc.save(output_path)
    return output_path


def pdf_to_word_hybrid_mode(pdf_path: str, output_path: Optional[str] = None, dpi: int = 200) -> str:
    """
    Hybrid mode: Renders pages as images for perfect layout, but also includes
    OCR/extracted text below each image for searchability and copy/paste.
    
    Best for complex documents where you need both visual fidelity AND text access.
    
    Args:
        pdf_path: Path to input PDF
        output_path: Optional output path
        dpi: Image resolution
    
    Returns:
        Path to output Word document
    """
    if output_path is None:
        output_path = str(Path(pdf_path).with_suffix('.docx'))
    
    pdf_doc = fitz.open(pdf_path)
    doc = DocxDocument()
    
    # Calculate zoom factor for desired DPI
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    
    for page_num in range(len(pdf_doc)):
        page = pdf_doc[page_num]
        
        # Add page break between pages (except first)
        if page_num > 0:
            doc.add_page_break()
        
        # Render page as high-quality image
        pix = page.get_pixmap(matrix=mat)
        
        # Save temp image
        temp_img = f"/tmp/pdf_page_{page_num}_{os.getpid()}.png"
        pix.save(temp_img)
        
        # Get image dimensions
        img_width_inches = page.rect.width / 72.0
        max_width = 6.5
        if img_width_inches > max_width:
            img_width_inches = max_width
        
        # Insert image into Word document
        try:
            doc.add_picture(temp_img, width=DocxInches(img_width_inches))
        except Exception as e:
            print(f"Error adding image for page {page_num}: {e}")
        
        # Cleanup temp file
        try:
            os.remove(temp_img)
        except:
            pass
        
        # Extract text from PDF (try direct extraction first, fall back to OCR)
        text = page.get_text().strip()
        
        # If no text extracted, try OCR
        if len(text) < 50:
            try:
                text = ocr_pdf_page(page, lang='eng+hin')
            except:
                text = ""
        
        # Add extracted/OCR text below the image (in smaller font, as reference)
        if text:
            doc.add_paragraph()  # Spacer
            para = doc.add_paragraph()
            para.add_run("--- Extracted Text (for copy/paste) ---").bold = True
            
            # Add the text content
            for line in text.split('\n'):
                if line.strip():
                    p = doc.add_paragraph(line.strip())
                    p.style = 'Normal'
    
    pdf_doc.close()
    doc.save(output_path)
    return output_path


def pdf_to_word_with_ocr(pdf_path: str, output_path: Optional[str] = None) -> str:
    """
    Convert PDF to Word with OCR support for image-based PDFs
    Automatically detects if OCR is needed
    
    Args:
        pdf_path: Path to input PDF
        output_path: Optional output path
    
    Returns:
        Path to output Word document
    """
    if output_path is None:
        output_path = str(Path(pdf_path).with_suffix('.docx'))
    
    # Check if PDF needs OCR
    needs_ocr = is_image_based_pdf(pdf_path)
    
    if needs_ocr:
        print("Image-based PDF detected, using OCR...")
        
        # Open PDF
        pdf_doc = fitz.open(pdf_path)
        doc = DocxDocument()
        
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            
            # Add page break between pages (except first)
            if page_num > 0:
                doc.add_page_break()
            
            # Perform OCR on page
            text = ocr_pdf_page(page)
            
            if text.strip():
                # Add OCR text to document
                for paragraph in text.split('\n\n'):
                    if paragraph.strip():
                        p = doc.add_paragraph(paragraph.strip())
                        p.style = 'Normal'
            
            # Also extract images
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = pdf_doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Save temp image
                    temp_img = f"/tmp/pdf_img_{page_num}_{img_index}.png"
                    with open(temp_img, "wb") as f:
                        f.write(image_bytes)
                    
                    # Insert into Word
                    doc.add_picture(temp_img, width=DocxInches(5))
                    
                    # Cleanup
                    os.remove(temp_img)
                except Exception:
                    pass
        
        pdf_doc.close()
        doc.save(output_path)
        return output_path
    else:
        # Use standard conversion for text-based PDFs
        return pdf_to_word_advanced(pdf_path, output_path)


def pdf_to_word_advanced(pdf_path: str, output_path: Optional[str] = None) -> str:
    """
    Convert PDF to Word with advanced font and formatting preservation
    Uses PyMuPDF for accurate text extraction with styles
    
    Args:
        pdf_path: Path to input PDF
        output_path: Optional output path
    
    Returns:
        Path to output Word document
    """
    if output_path is None:
        output_path = str(Path(pdf_path).with_suffix('.docx'))
    
    # Open PDF with PyMuPDF for accurate extraction
    pdf_doc = fitz.open(pdf_path)
    
    # Create Word document
    doc = DocxDocument()
    
    for page_num in range(len(pdf_doc)):
        page = pdf_doc[page_num]
        
        # Extract text blocks with formatting
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
        
        # Add page break between pages (except first)
        if page_num > 0:
            doc.add_page_break()
        
        for block in blocks:
            if block["type"] == 0:  # Text block
                for line in block["lines"]:
                    # Create paragraph for each line
                    para = doc.add_paragraph()
                    
                    for span in line["spans"]:
                        text = span["text"]
                        if not text.strip():
                            continue
                        
                        # Create run with formatting
                        run = para.add_run(text)
                        
                        # Apply font size
                        run.font.size = DocxPt(span["size"])
                        
                        # Apply font name (map common PDF fonts to Word fonts)
                        pdf_font = span["font"].lower()
                        if "times" in pdf_font:
                            run.font.name = "Times New Roman"
                        elif "arial" in pdf_font or "helvetica" in pdf_font:
                            run.font.name = "Arial"
                        elif "courier" in pdf_font:
                            run.font.name = "Courier New"
                        elif "georgia" in pdf_font:
                            run.font.name = "Georgia"
                        elif "verdana" in pdf_font:
                            run.font.name = "Verdana"
                        else:
                            run.font.name = "Calibri"  # Default
                        
                        # Apply bold/italic from flags
                        flags = span["flags"]
                        run.bold = bool(flags & 2**4)  # Bold flag
                        run.italic = bool(flags & 2**1)  # Italic flag
                        
                        # Apply color
                        color = span["color"]
                        if color != 0:
                            r = (color >> 16) & 0xFF
                            g = (color >> 8) & 0xFF
                            b = color & 0xFF
                            run.font.color.rgb = RGBColor(r, g, b)
            
            elif block["type"] == 1:  # Image block
                # Extract and insert image
                try:
                    img_rect = fitz.Rect(block["bbox"])
                    clip = page.get_pixmap(clip=img_rect, matrix=fitz.Matrix(2, 2))
                    img_bytes = clip.tobytes("png")
                    
                    # Save temp image
                    temp_img = f"/tmp/pdf_img_{page_num}_{block['number']}.png"
                    with open(temp_img, "wb") as f:
                        f.write(img_bytes)
                    
                    # Insert into Word
                    doc.add_picture(temp_img, width=DocxInches(5))
                    os.remove(temp_img)
                except Exception:
                    pass
    
    pdf_doc.close()
    doc.save(output_path)
    return output_path


def pdf_to_word(pdf_path: str, output_path: Optional[str] = None, mode: str = 'auto') -> str:
    """
    Convert PDF to Word document with smart auto-detection.
    
    Args:
        pdf_path: Path to input PDF
        output_path: Optional output path
        mode: Conversion mode
            - 'auto' (default): Smart auto-detection based on PDF analysis
            - 'hybrid': Image + extracted text (best for complex docs like Aadhaar)
            - 'image': Image only (perfect visual layout, no text)
            - 'text': Extract editable text (may break layout for complex PDFs)
            - 'ocr': Force OCR (for scanned documents)
    
    Returns:
        Path to output Word document
    """
    if output_path is None:
        output_path = str(Path(pdf_path).with_suffix('.docx'))
    
    # Smart auto-detection
    if mode == 'auto':
        analysis = analyze_pdf_type(pdf_path)
        print(f"[PDF Analysis] Type: {analysis['recommendation']}, "
              f"Text: {analysis['text_length']} chars, "
              f"Images: {analysis['image_count']}, "
              f"Complex: {analysis['is_complex']}")
        
        recommendation = analysis['recommendation']
        
        if recommendation == 'text':
            # Simple text PDF - try to extract editable text
            try:
                cv = PDFToDocxConverter(pdf_path)
                cv.convert(output_path)
                cv.close()
                return output_path
            except Exception:
                # Fallback to advanced method
                return pdf_to_word_advanced(pdf_path, output_path)
        
        elif recommendation == 'ocr':
            # Scanned document - use OCR
            return pdf_to_word_with_ocr(pdf_path, output_path)
        
        elif recommendation == 'hybrid':
            # Complex layout with images - use hybrid mode (image + text)
            return pdf_to_word_hybrid_mode(pdf_path, output_path)
        
        else:  # 'image' or unknown
            # Default to image mode for safety
            return pdf_to_word_image_mode(pdf_path, output_path)
    
    # Explicit mode selection
    elif mode == 'hybrid':
        return pdf_to_word_hybrid_mode(pdf_path, output_path)
    
    elif mode == 'image':
        return pdf_to_word_image_mode(pdf_path, output_path)
    
    elif mode == 'ocr':
        return pdf_to_word_with_ocr(pdf_path, output_path)
    
    elif mode == 'text':
        try:
            cv = PDFToDocxConverter(pdf_path)
            cv.convert(output_path)
            cv.close()
            return output_path
        except Exception:
            return pdf_to_word_advanced(pdf_path, output_path)
    
    else:
        # Unknown mode - default to hybrid for safety
        return pdf_to_word_hybrid_mode(pdf_path, output_path)


def pdf_to_excel_with_ocr(pdf_path: str, output_path: Optional[str] = None) -> str:
    """
    Extract data from image-based PDF to Excel using OCR
    """
    if output_path is None:
        output_path = str(Path(pdf_path).with_suffix('.xlsx'))
    
    pdf_doc = fitz.open(pdf_path)
    workbook = Workbook()
    ws = workbook.active
    ws.title = "PDF Data"
    
    row_num = 1
    
    for page_num in range(len(pdf_doc)):
        page = pdf_doc[page_num]
        
        # Add page header
        ws.cell(row=row_num, column=1, value=f"Page {page_num + 1}")
        ws.cell(row=row_num, column=1).font = Font(bold=True, size=12)
        row_num += 1
        
        # Perform OCR
        text = ocr_pdf_page(page)
        
        if text.strip():
            # Split by lines and add to Excel
            for line in text.split('\n'):
                if line.strip():
                    # Try to detect columns (split by multiple spaces or tabs)
                    parts = [p.strip() for p in line.split('  ') if p.strip()]
                    
                    if len(parts) > 1:
                        # Multiple columns detected
                        for col_idx, part in enumerate(parts):
                            ws.cell(row=row_num, column=col_idx + 1, value=part)
                    else:
                        # Single column
                        ws.cell(row=row_num, column=1, value=line.strip())
                    
                    row_num += 1
        
        row_num += 1  # Space between pages
    
    # Auto-adjust column widths
    for column in ws.columns:
        max_length = 0
        column_letter = column[0].column_letter
        for cell in column:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        adjusted_width = min(max_length + 2, 50)
        ws.column_dimensions[column_letter].width = adjusted_width
    
    pdf_doc.close()
    workbook.save(output_path)
    return output_path


def pdf_to_excel_advanced(pdf_path: str, output_path: Optional[str] = None) -> str:
    """
    Extract tables from PDF to Excel with formatting preservation
    Uses PyMuPDF for accurate table detection
    """
    if output_path is None:
        output_path = str(Path(pdf_path).with_suffix('.xlsx'))
    
    pdf_doc = fitz.open(pdf_path)
    workbook = Workbook()
    ws = workbook.active
    ws.title = "PDF Data"
    
    # Style for headers
    header_font = Font(bold=True, size=11)
    header_fill = PatternFill(start_color="DAEEF3", end_color="DAEEF3", fill_type="solid")
    thin_border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    row_num = 1
    
    for page_num in range(len(pdf_doc)):
        page = pdf_doc[page_num]
        
        # Try to find tables using PyMuPDF's table detection
        tables = page.find_tables()
        
        if tables:
            for table in tables:
                # Add page header
                ws.cell(row=row_num, column=1, value=f"Page {page_num + 1}")
                ws.cell(row=row_num, column=1).font = Font(bold=True, size=12)
                row_num += 1
                
                # Extract table data
                for row_idx, row in enumerate(table.extract()):
                    for col_idx, cell in enumerate(row):
                        cell_obj = ws.cell(row=row_num, column=col_idx + 1, value=cell or "")
                        cell_obj.border = thin_border
                        
                        # Style first row as header
                        if row_idx == 0:
                            cell_obj.font = header_font
                            cell_obj.fill = header_fill
                    
                    row_num += 1
                
                row_num += 1  # Space between tables
        else:
            # No tables found - extract text in columns
            text = page.get_text("text")
            ws.cell(row=row_num, column=1, value=f"--- Page {page_num + 1} ---")
            ws.cell(row=row_num, column=1).font = Font(bold=True)
            row_num += 1
            
            for line in text.split('\n'):
                if line.strip():
                    ws.cell(row=row_num, column=1, value=line.strip())
                    row_num += 1
            
            row_num += 1
    
    # Auto-adjust column widths
    for column in ws.columns:
        max_length = 0
        column_letter = column[0].column_letter
        for cell in column:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        adjusted_width = min(max_length + 2, 50)
        ws.column_dimensions[column_letter].width = adjusted_width
    
    pdf_doc.close()
    workbook.save(output_path)
    return output_path


def pdf_to_excel(pdf_path: str, output_path: Optional[str] = None) -> str:
    """
    Extract tables from PDF and convert to Excel
    Automatically uses OCR for image-based PDFs
    """
    if output_path is None:
        output_path = str(Path(pdf_path).with_suffix('.xlsx'))
    
    # Check if OCR is needed
    if is_image_based_pdf(pdf_path):
        return pdf_to_excel_with_ocr(pdf_path, output_path)
    
    # Try camelot first if available (best for visible table borders)
    if CAMELOT_AVAILABLE:
        try:
            tables = camelot.read_pdf(pdf_path, pages='all', flavor='lattice')
            if len(tables) > 0:
                workbook = Workbook()
                default_sheet = workbook.active
                
                for i, table in enumerate(tables):
                    sheet = workbook.create_sheet(title=f"Table_{i+1}")
                    df = table.df
                    for r_idx, row in enumerate(df.values, 1):
                        for c_idx, value in enumerate(row, 1):
                            sheet.cell(row=r_idx, column=c_idx, value=str(value))
                
                if len(workbook.sheetnames) > 1:
                    workbook.remove(default_sheet)
                
                workbook.save(output_path)
                return output_path
        except Exception:
            pass
    
    # Use PyMuPDF advanced extraction
    return pdf_to_excel_advanced(pdf_path, output_path)


def pdf_to_pptx_with_ocr(pdf_path: str, output_path: Optional[str] = None) -> str:
    """
    Convert image-based PDF to PowerPoint with OCR
    """
    if output_path is None:
        output_path = str(Path(pdf_path).with_suffix('.pptx'))
    
    pdf_doc = fitz.open(pdf_path)
    prs = Presentation()
    
    # Set slide size
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Blank layout
    blank_layout = prs.slide_layouts[6]
    
    for page_num in range(len(pdf_doc)):
        page = pdf_doc[page_num]
        
        # Add slide
        slide = prs.slides.add_slide(blank_layout)
        
        # Render page as image
        zoom = 2.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        
        # Save temp image
        temp_img_path = f"/tmp/pdf_slide_{page_num}.png"
        pix.save(temp_img_path)
        
        # Add image to slide
        margin = Inches(0.2)
        pic = slide.shapes.add_picture(
            temp_img_path,
            margin,
            margin,
            width=prs.slide_width - (margin * 2),
            height=prs.slide_height - (margin * 2)
        )
        
        # Perform OCR and add as text box
        text = ocr_pdf_page(page)
        if text.strip():
            # Add text box at bottom
            left = Inches(0.5)
            top = prs.slide_height - Inches(2)
            width = prs.slide_width - Inches(1)
            height = Inches(1.5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = text[:500]  # Limit text length
            text_frame.word_wrap = True
            
            # Style text
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(10)
        
        # Cleanup
        os.remove(temp_img_path)
    
    pdf_doc.close()
    prs.save(output_path)
    return output_path


def pdf_to_pptx_advanced(pdf_path: str, output_path: Optional[str] = None) -> str:
    """
    Convert PDF to PowerPoint with high-quality rendering
    Preserves text as editable where possible
    """
    if output_path is None:
        output_path = str(Path(pdf_path).with_suffix('.pptx'))
    
    pdf_doc = fitz.open(pdf_path)
    
    # Create presentation with proper aspect ratio
    prs = Presentation()
    
    # Get first page dimensions for slide size
    first_page = pdf_doc[0]
    page_rect = first_page.rect
    
    # Set slide size based on PDF proportions
    aspect_ratio = page_rect.width / page_rect.height
    if aspect_ratio > 1.5:
        # Wide format (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        # Standard format (4:3)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    # Blank layout
    blank_layout = prs.slide_layouts[6]
    
    for page_num in range(len(pdf_doc)):
        page = pdf_doc[page_num]
        
        # Render page as high-quality image
        zoom = 2.0  # 2x for better quality
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        
        # Save temp image
        temp_img_path = f"/tmp/pdf_slide_{page_num}.png"
        pix.save(temp_img_path)
        
        # Add slide
        slide = prs.slides.add_slide(blank_layout)
        
        # Calculate image size to fit slide
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Add image centered on slide
        margin = Inches(0.2)
        pic = slide.shapes.add_picture(
            temp_img_path,
            margin,
            margin,
            width=slide_width - (margin * 2),
            height=slide_height - (margin * 2)
        )
        
        # Cleanup temp image
        os.remove(temp_img_path)
    
    pdf_doc.close()
    prs.save(output_path)
    return output_path


def pdf_to_pptx(pdf_path: str, output_path: Optional[str] = None) -> str:
    """
    Convert PDF to PowerPoint presentation
    Automatically uses OCR for image-based PDFs
    """
    if output_path is None:
        output_path = str(Path(pdf_path).with_suffix('.pptx'))
    
    # Check if OCR is needed
    if is_image_based_pdf(pdf_path):
        return pdf_to_pptx_with_ocr(pdf_path, output_path)
    
    return pdf_to_pptx_advanced(pdf_path, output_path)


def word_to_pdf(docx_path: str, output_path: Optional[str] = None) -> str:
    """
    Convert Word document to PDF
    Uses reportlab for PDF generation
    """
    if output_path is None:
        output_path = str(Path(docx_path).with_suffix('.pdf'))
    
    doc = DocxDocument(docx_path)
    
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    
    y_position = height - inch
    line_height = 14
    
    for para in doc.paragraphs:
        text = para.text
        if not text.strip():
            y_position -= line_height
            continue
        
        if y_position < inch:
            c.showPage()
            y_position = height - inch
        
        # Handle long lines with word wrap
        words = text.split()
        current_line = ""
        
        for word in words:
            test_line = current_line + " " + word if current_line else word
            if len(test_line) > 80:
                c.drawString(inch, y_position, current_line)
                y_position -= line_height
                current_line = word
                if y_position < inch:
                    c.showPage()
                    y_position = height - inch
            else:
                current_line = test_line
        
        if current_line:
            c.drawString(inch, y_position, current_line)
            y_position -= line_height
    
    c.save()
    return output_path


def excel_to_pdf(xlsx_path: str, output_path: Optional[str] = None) -> str:
    """Convert Excel to PDF"""
    from openpyxl import load_workbook
    
    if output_path is None:
        output_path = str(Path(xlsx_path).with_suffix('.pdf'))
    
    wb = load_workbook(xlsx_path)
    
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    
    for sheet in wb.worksheets:
        y_position = height - inch
        c.setFont("Helvetica-Bold", 14)
        c.drawString(inch, y_position, f"Sheet: {sheet.title}")
        y_position -= 25
        c.setFont("Helvetica", 10)
        
        for row in sheet.iter_rows(values_only=True):
            if y_position < inch:
                c.showPage()
                y_position = height - inch
            
            row_text = " | ".join(str(cell) if cell else "" for cell in row)
            if len(row_text) > 100:
                row_text = row_text[:100] + "..."
            
            c.drawString(inch, y_position, row_text)
            y_position -= 14
        
        c.showPage()
    
    c.save()
    return output_path


def pptx_to_pdf(pptx_path: str, output_path: Optional[str] = None) -> str:
    """Convert PowerPoint to PDF"""
    if output_path is None:
        output_path = str(Path(pptx_path).with_suffix('.pdf'))
    
    prs = Presentation(pptx_path)
    
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    
    for i, slide in enumerate(prs.slides):
        y_position = height - inch
        c.setFont("Helvetica-Bold", 16)
        c.drawString(inch, y_position, f"Slide {i + 1}")
        y_position -= 30
        c.setFont("Helvetica", 12)
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text
                for line in text.split('\n'):
                    if y_position < inch:
                        c.showPage()
                        y_position = height - inch
                    
                    if len(line) > 80:
                        line = line[:80] + "..."
                    
                    c.drawString(inch, y_position, line)
                    y_position -= 14
        
        c.showPage()
    
    c.save()
    return output_path


def image_to_pdf(image_paths: List[str], output_path: str) -> str:
    """Convert images to PDF"""
    from PIL import Image
    
    images = []
    for path in image_paths:
        img = Image.open(path)
        if img.mode == 'RGBA':
            img = img.convert('RGB')
        images.append(img)
    
    if images:
        images[0].save(
            output_path,
            save_all=True,
            append_images=images[1:] if len(images) > 1 else [],
            format='PDF',
            resolution=150
        )
    
    return output_path


def html_to_pdf(html_content: str, output_path: str) -> str:
    """Convert HTML to PDF"""
    from reportlab.platypus import SimpleDocTemplate, Paragraph
    from reportlab.lib.styles import getSampleStyleSheet
    import re
    
    clean_text = re.sub(r'<[^>]+>', '', html_content)
    
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    
    story = []
    for para in clean_text.split('\n\n'):
        if para.strip():
            story.append(Paragraph(para, styles['Normal']))
    
    if story:
        doc.build(story)
    else:
        c = canvas.Canvas(output_path, pagesize=letter)
        c.save()
    
    return output_path
